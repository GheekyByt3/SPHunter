"""
SPHunter Detection Module

Analyzes discovered files for sensitive content using two layers:
1. Filename pattern matching (fast, no download)
2. Content regex scanning (requires downloaded file)
"""

import os
import re
import yaml
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeout
from rich.console import Console

console = Console()

CONTENT_TEXT_CAP = 50 * 1024   # 50KB max text per file
CONTENT_PARSE_TIMEOUT = 10      # seconds before giving up on a file

DEFAULT_RULES_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "config", "rules.yaml")

SEVERITY_SCORES = {
    "black": 4,
    "red": 3,
    "yellow": 2,
    "green": 1,
}

SEVERITY_COLORS = {
    "black": "bold white on black",
    "red": "bold red",
    "yellow": "yellow",
    "green": "green",
}

# Map old severity names to new Snaffler-style names
SEVERITY_MAP = {
    "critical": "black",
    "high": "red",
    "medium": "yellow",
    "low": "green",
    "info": "green",
    "black": "black",
    "red": "red",
    "yellow": "yellow",
    "green": "green",
}


class SensitiveFileDetector:
    """Detects sensitive files by filename and content pattern matching."""

    def __init__(self, rules_path: str = None):
        self.rules_path = rules_path or DEFAULT_RULES_PATH
        self.filename_rules = []
        self.content_rules = []
        self.findings = []
        self._load_rules()

    def _load_rules(self):
        """Load detection rules from YAML config."""
        try:
            with open(self.rules_path, "r") as f:
                rules = yaml.safe_load(f)

            self.filename_rules = rules.get("filename_rules", [])
            self.content_rules = rules.get("content_rules", [])

            # Pre-compile regex patterns
            for rule in self.filename_rules:
                rule["_compiled"] = re.compile(rule["pattern"], re.IGNORECASE)

            for rule in self.content_rules:
                rule["_compiled"] = re.compile(rule["pattern"])

            console.print(f"[green][+] Loaded {len(self.filename_rules)} filename rules, {len(self.content_rules)} content rules[/green]")

        except FileNotFoundError:
            console.print(f"[red][-] Rules file not found: {self.rules_path}[/red]")
            console.print("[yellow][*] Running with empty ruleset[/yellow]")
        except yaml.YAMLError as e:
            console.print(f"[red][-] Error parsing rules YAML: {e}[/red]")

    def analyze_files(self, files: list, content_inspection: bool = True) -> list:
        """
        Run detection against all discovered files.

        Args:
            files: List of file metadata dicts from the crawler.
            content_inspection: Whether to scan file contents.

        Returns:
            List of finding dicts.
        """
        console.print(f"\n[yellow][*] Analyzing {len(files)} files for sensitive data...[/yellow]")

        for file_info in files:
            file_findings = []

            # Layer 1: Filename matching
            filename_hits = self._check_filename(file_info["name"])
            file_findings.extend(filename_hits)

            # Layer 2: Content inspection — skip if filename already hit black
            filename_max_sev = max(
                (SEVERITY_SCORES.get(h["severity"], 0) for h in filename_hits), default=0
            )
            if content_inspection and file_info.get("local_path") and filename_max_sev < SEVERITY_SCORES["black"]:
                content_hits = self._check_content(file_info["local_path"], file_info["name"])
                file_findings.extend(content_hits)

            if file_findings:
                # Attach findings to file info and add to master list
                file_info["findings"] = file_findings
                highest_severity = max(file_findings, key=lambda f: SEVERITY_SCORES.get(f["severity"], 0))

                finding_record = {
                    "file": file_info,
                    "findings": file_findings,
                    "highest_severity": highest_severity["severity"],
                    "severity_score": SEVERITY_SCORES.get(highest_severity["severity"], 0),
                }
                self.findings.append(finding_record)

                # Live console output
                sev = highest_severity["severity"]
                color = SEVERITY_COLORS.get(sev, "white")
                rules_hit = ", ".join(f["rule_name"] for f in file_findings)
                console.print(
                    f"  [{color}][{sev.upper()}][/{color}] "
                    f"{file_info['siteName']} / {file_info['fullPath']} "
                    f"[dim]({rules_hit})[/dim]"
                )

        # Sort findings by severity
        self.findings.sort(key=lambda f: f["severity_score"], reverse=True)

        # Summary
        severity_counts = {}
        for f in self.findings:
            sev = f["highest_severity"]
            severity_counts[sev] = severity_counts.get(sev, 0) + 1

        console.print(f"\n[green][+] Detection complete: {len(self.findings)} files flagged[/green]")
        for sev in ["black", "red", "yellow", "green"]:
            count = severity_counts.get(sev, 0)
            if count > 0:
                color = SEVERITY_COLORS.get(sev, "white")
                console.print(f"    [{color}]{sev.upper()}: {count}[/{color}]")

        return self.findings

    def _check_filename(self, filename: str) -> list:
        """Check filename against all filename rules."""
        hits = []
        for rule in self.filename_rules:
            if rule["_compiled"].search(filename):
                hits.append({
                    "rule_name": rule["name"],
                    "severity": SEVERITY_MAP.get(rule["severity"], rule["severity"]),
                    "description": rule["description"],
                    "match_type": "filename",
                    "matched_value": filename,
                })
        return hits

    def _check_content(self, local_path: str, filename: str) -> list:
        """Scan file content against all content rules."""
        hits = []
        content = self._extract_text(local_path, filename)
        if not content:
            return hits

        for rule in self.content_rules:
            matches = rule["_compiled"].findall(content)
            if matches:
                sample = matches[0] if isinstance(matches[0], str) else str(matches[0])
                hits.append({
                    "rule_name": rule["name"],
                    "severity": SEVERITY_MAP.get(rule["severity"], rule["severity"]),
                    "description": rule["description"],
                    "match_type": "content",
                    "matched_value": self._mask_value(sample),
                    "match_count": len(matches),
                })

        return hits

    def _extract_text(self, local_path: str, filename: str) -> str:
        """Extract readable text from a file, dispatched by extension with timeout + size cap."""
        try:
            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(self._parse_file, local_path, filename)
                text = future.result(timeout=CONTENT_PARSE_TIMEOUT)
                return (text or "")[:CONTENT_TEXT_CAP]
        except FuturesTimeout:
            console.print(f"[dim]    [!] Parse timeout on {filename}, skipping content[/dim]")
            return ""
        except Exception as e:
            console.print(f"[dim]    [!] Could not parse {filename}: {e}[/dim]")
            return ""

    @staticmethod
    def _parse_file(local_path: str, filename: str) -> str:
        """Dispatch to the right parser based on file extension."""
        ext = os.path.splitext(filename.lower())[1]

        if ext == ".docx":
            import docx
            doc = docx.Document(local_path)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(local_path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    parts.append(" ".join(str(c) for c in row if c is not None))
            wb.close()
            return "\n".join(parts)

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(local_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)

        if ext == ".pdf":
            from pdfminer.high_level import extract_text as pdf_extract
            return pdf_extract(local_path) or ""

        if ext == ".msg":
            import extract_msg
            with extract_msg.Message(local_path) as msg:
                return f"{msg.subject or ''}\n{msg.body or ''}"

        if ext == ".rtf":
            from striprtf.striprtf import rtf_to_text
            with open(local_path, "r", errors="ignore") as f:
                return rtf_to_text(f.read())

        with open(local_path, "r", errors="ignore", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def _mask_value(value: str) -> str:
        """Mask sensitive values for safe reporting."""
        value = value.strip()
        if len(value) <= 8:
            return value[:2] + "*" * (len(value) - 2)
        return value[:4] + "*" * (len(value) - 8) + value[-4:]
